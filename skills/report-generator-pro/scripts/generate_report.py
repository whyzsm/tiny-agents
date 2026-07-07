#!/usr/bin/env python3
"""
generate_report.py - 生成报告文件（Word/Excel/PPT/PDF/Markdown）
"""
import argparse
import sys
import os
import json
from datetime import datetime, date

def ensure_dir(path):
    os.makedirs(os.path.dirname(path), exist_ok=True)

def generate_markdown(content_dict, report_type, target_date, author_info):
    """生成 Markdown 格式报告"""
    today = target_date or date.today().strftime('%Y年%m月%d日')
    name = author_info.get('name', '')
    dept = author_info.get('department', '')
    title = f"# {name}{dept} {today} {report_type}\n\n"

    sections = []
    if content_dict.get('highlights'):
        sections.append(f"## 📊 核心亮点\n{content_dict['highlights']}\n")
    if content_dict.get('completed_tasks'):
        sections.append(f"## ✅ 已完成工作\n{content_dict['completed_tasks']}\n")
    if content_dict.get('in_progress_tasks'):
        sections.append(f"## 🔄 进行中工作\n{content_dict['in_progress_tasks']}\n")
    if content_dict.get('key_metrics'):
        sections.append(f"## 📈 关键数据\n{content_dict['key_metrics']}\n")
    if content_dict.get('blocked_items'):
        sections.append(f"## ⚠️ 待解决问题\n{content_dict['blocked_items']}\n")
    if content_dict.get('next_plans'):
        sections.append(f"## 📋 下阶段计划\n{content_dict['next_plans']}\n")

    footer = f"\n---\n*报告生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}*\n"
    return title + "\n".join(sections) + footer

def generate_word(content_dict, report_type, target_date, author_info, output_path):
    """生成 Word 文档"""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        today = target_date or date.today().strftime('%Y年%m月%d日')
        name = author_info.get('name', '')
        dept = author_info.get('department', '')

        # 标题
        title = doc.add_heading(f'{name}{dept} {today} {report_type}', level=1)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 各章节
        section_map = [
            ('highlights', '核心亮点', '📊'),
            ('completed_tasks', '已完成工作', '✅'),
            ('in_progress_tasks', '进行中工作', '🔄'),
            ('key_metrics', '关键数据', '📈'),
            ('blocked_items', '待解决问题', '⚠️'),
            ('next_plans', '下阶段计划', '📋'),
        ]

        for key, label, icon in section_map:
            if content_dict.get(key):
                doc.add_heading(f'{icon} {label}', level=2)
                content = content_dict[key]
                if isinstance(content, list):
                    for item in content:
                        doc.add_paragraph(item, style='List Bullet')
                else:
                    for line in str(content).split('\n'):
                        if line.strip():
                            doc.add_paragraph(line.strip())

        # 页脚
        doc.add_paragraph(f'报告生成时间：{datetime.now().strftime("%Y-%m-%d %H:%M")}')

        ensure_dir(output_path)
        doc.save(output_path)
        return {"success": True, "path": output_path}
    except ImportError:
        return {"error": "需要安装 python-docx: pip install python-docx"}
    except Exception as e:
        return {"error": f"生成 Word 失败: {e}"}

def generate_excel(content_dict, report_type, target_date, author_info, output_path):
    """生成 Excel 报告"""
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        wb = openpyxl.Workbook()
        ws = wb.active
        today = target_date or date.today().strftime('%Y-%m-%d')
        ws.title = f'{report_type}_{today}'

        # 标题行
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        ws['A1'] = f'{author_info.get("name", "")} {report_type} - {today}'
        ws['A1'].font = Font(bold=True, size=14, color="FFFFFF")
        ws['A1'].fill = header_fill
        ws.merge_cells('A1:C1')

        row = 3
        section_map = [
            ('completed_tasks', '已完成工作'),
            ('in_progress_tasks', '进行中工作'),
            ('key_metrics', '关键数据'),
            ('blocked_items', '待解决问题'),
            ('next_plans', '下阶段计划'),
        ]

        for key, label in section_map:
            if content_dict.get(key):
                ws.cell(row=row, column=1, value=label).font = Font(bold=True, size=12)
                row += 1
                content = content_dict[key]
                items = content if isinstance(content, list) else str(content).split('\n')
                for item in items:
                    if item.strip():
                        ws.cell(row=row, column=1, value=item.strip())
                        row += 1
                row += 1

        ws.column_dimensions['A'].width = 60
        ensure_dir(output_path)
        wb.save(output_path)
        return {"success": True, "path": output_path}
    except ImportError:
        return {"error": "需要安装 openpyxl: pip install openpyxl"}
    except Exception as e:
        return {"error": f"生成 Excel 失败: {e}"}

def generate_ppt(content_dict, report_type, target_date, author_info, output_path):
    """生成 PPT 报告"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        today = target_date or date.today().strftime('%Y年%m月%d日')
        name = author_info.get('name', '')

        # 封面
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f'{report_type}'
        slide.placeholders[1].text = f'{name}  {today}'

        # 内容页
        section_map = [
            ('completed_tasks', '已完成工作 ✅'),
            ('in_progress_tasks', '进行中工作 🔄'),
            ('key_metrics', '关键数据 📈'),
            ('blocked_items', '待解决问题 ⚠️'),
            ('next_plans', '下阶段计划 📋'),
        ]

        for key, label in section_map:
            if content_dict.get(key):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = label
                content = content_dict[key]
                items = content if isinstance(content, list) else str(content).split('\n')
                body = '\n'.join([f'• {i.strip()}' for i in items if i.strip()])
                slide.placeholders[1].text = body

        ensure_dir(output_path)
        prs.save(output_path)
        return {"success": True, "path": output_path}
    except ImportError:
        return {"error": "需要安装 python-pptx: pip install python-pptx"}
    except Exception as e:
        return {"error": f"生成 PPT 失败: {e}"}

def generate_pdf(content_dict, report_type, target_date, author_info, output_path):
    """生成 PDF 报告"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        today = target_date or date.today().strftime('%Y年%m月%d日')
        name = author_info.get('name', '')
        dept = author_info.get('department', '')

        ensure_dir(output_path)
        doc = SimpleDocTemplate(output_path, pagesize=A4, topMargin=20*mm, bottomMargin=20*mm)
        story = []
        styles = getSampleStyleSheet()

        # 尝试注册中文字体
        chinese_font = 'Helvetica'
        for font_name in ['SimSun', 'Microsoft YaHei', 'SimHei', 'STSong']:
            try:
                pdfmetrics.registerFont(TTFont(font_name, f'C:/Windows/Fonts/{font_name}.ttf'))
                chinese_font = font_name
                break
            except Exception:
                continue

        title_style = ParagraphStyle('ChTitle', parent=styles['Title'], fontName=chinese_font, fontSize=18, alignment=TA_CENTER, spaceAfter=10)
        heading_style = ParagraphStyle('ChHeading', parent=styles['Heading2'], fontName=chinese_font, fontSize=14, spaceBefore=12, spaceAfter=6)
        body_style = ParagraphStyle('ChBody', parent=styles['Normal'], fontName=chinese_font, fontSize=11, leading=18)

        story.append(Paragraph(f'{name}{dept} {today} {report_type}', title_style))
        story.append(HRFlowable(width="100%", thickness=1, color="#1E88E5"))
        story.append(Spacer(1, 8*mm))

        section_map = [
            ('highlights', '核心亮点'),
            ('completed_tasks', '已完成工作'),
            ('in_progress_tasks', '进行中工作'),
            ('key_metrics', '关键数据'),
            ('blocked_items', '待解决问题'),
            ('next_plans', '下阶段计划'),
        ]

        for key, label in section_map:
            if content_dict.get(key):
                story.append(Paragraph(f'■ {label}', heading_style))
                content = content_dict[key]
                items = content if isinstance(content, list) else str(content).split('\n')
                for item in items:
                    if item.strip():
                        story.append(Paragraph(f'• {item.strip()}', body_style))
                story.append(Spacer(1, 4*mm))

        story.append(Spacer(1, 6*mm))
        story.append(HRFlowable(width="100%", thickness=0.5, color="#CCCCCC"))
        story.append(Paragraph(f'报告生成时间：{datetime.now().strftime("%Y-%m-%d %H:%M")}',
                               ParagraphStyle('Footer', parent=styles['Normal'], fontName=chinese_font, fontSize=9, textColor='#999999')))

        doc.build(story)
        return {"success": True, "path": output_path}
    except ImportError:
        return {"error": "需要安装 reportlab: pip install reportlab"}
    except Exception as e:
        return {"error": f"生成 PDF 失败: {e}"}

def main():
    parser = argparse.ArgumentParser(description='生成报告文件')
    parser.add_argument('--type', required=True, help='报告类型（日报/周报/月报/季报/年报）')
    parser.add_argument('--format', required=True, help='输出格式（word/excel/ppt/pdf/markdown）')
    parser.add_argument('--content', required=True, help='内容 JSON 字符串或 JSON 文件路径')
    parser.add_argument('--output', required=True, help='输出文件路径')
    parser.add_argument('--date', help='报告日期（YYYY年MM月DD日格式）')
    parser.add_argument('--author', default='{}', help='作者信息 JSON（name/department）')
    args = parser.parse_args()

    # 解析内容
    if os.path.exists(args.content):
        with open(args.content, 'r', encoding='utf-8') as f:
            content_dict = json.load(f)
    else:
        content_dict = json.loads(args.content)

    author_info = json.loads(args.author)

    fmt = args.format.lower()
    if fmt in ('word', 'docx'):
        result = generate_word(content_dict, args.type, args.date, author_info, args.output)
    elif fmt in ('excel', 'xlsx'):
        result = generate_excel(content_dict, args.type, args.date, author_info, args.output)
    elif fmt in ('ppt', 'pptx'):
        result = generate_ppt(content_dict, args.type, args.date, author_info, args.output)
    elif fmt == 'markdown':
        md = generate_markdown(content_dict, args.type, args.date, author_info)
        ensure_dir(args.output)
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(md)
        result = {"success": True, "path": args.output}
    elif fmt == 'pdf':
        result = generate_pdf(content_dict, args.type, args.date, author_info, args.output)
    else:
        result = {"error": f"不支持的格式: {fmt}，支持: word/excel/ppt/pdf/markdown"}

    print(json.dumps(result, ensure_ascii=False, indent=2))

if __name__ == '__main__':
    main()
