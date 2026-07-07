#!/usr/bin/env python3
"""
doc_image_extractor.py

从多种文档格式中提取内嵌图片。
支持格式：PPTX、DOCX、XLSX

依赖安装：
    pip install PyMuPDF python-pptx python-docx openpyxl

用法：
    # 单个文件
    python doc_image_extractor.py <文档路径> [输出目录]

    # 多个文件
    python doc_image_extractor.py file1.pdf file2.pptx file3.docx [输出目录]

    # 批量扫描目录
    python doc_image_extractor.py --dir <目录路径> [输出目录]
    python doc_image_extractor.py -d <目录路径> [输出目录]
"""

import sys
import os
import argparse
from pathlib import Path
from typing import Optional, List, Dict

MIN_IMAGE_SIZE = 100 * 1024

# -------------------- 依赖检查 --------------------

def _require(package_name: str, install_cmd: str):
    """尝试导入包，失败时提示安装并退出。"""
    try:
        __import__(package_name)
    except ImportError:
        print(f"❌ 缺少依赖: {package_name}")
        print(f"   请运行: {install_cmd}")
        sys.exit(1)


# -------------------- PPTX 图片提取 --------------------

def extract_images_from_pptx(filepath: str, output_dir: str) -> int:
    """使用 python-pptx 从 PPTX 中提取图片。"""
    _require("pptx", "pip install python-pptx")
    from pptx import Presentation

    prs = Presentation(filepath)
    image_count = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "image") and shape.image:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    # 跳过过小图片
                    if len(image_bytes) < MIN_IMAGE_SIZE:
                        continue

                    # 确定扩展名
                    content_type = image.content_type
                    ext = content_type.split("/")[-1] if "/" in content_type else "png"
                    if ext == "jpeg":
                        ext = "jpg"

                    image_count += 1
                    image_filename = f"pptx_s{slide_num}_img{image_count}.{ext}"
                    image_path = os.path.join(output_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                except Exception:
                    continue

    return image_count


# -------------------- DOCX 图片提取 --------------------

def extract_images_from_docx(filepath: str, output_dir: str) -> int:
    """使用 python-docx 从 DOCX 中提取图片。"""
    _require("docx", "pip install python-docx")
    from docx import Document

    doc = Document(filepath)
    image_count = 0

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                image = rel.target_part
                image_bytes = image.blob
                # 跳过过小图片
                if len(image_bytes) < MIN_IMAGE_SIZE:
                    continue

                # 确定扩展名
                content_type = image.content_type if hasattr(image, "content_type") else "image/png"
                ext = content_type.split("/")[-1] if "/" in content_type else "png"
                if ext == "jpeg":
                    ext = "jpg"

                image_count += 1
                image_filename = f"docx_img{image_count}.{ext}"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
            except Exception:
                continue

    return image_count


# -------------------- XLSX 图片提取 --------------------

def extract_images_from_xlsx(filepath: str, output_dir: str) -> int:
    """使用 openpyxl 从 XLSX 中提取图片。"""
    _require("openpyxl", "pip install openpyxl")
    from openpyxl import load_workbook

    wb = load_workbook(filepath)
    image_count = 0

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        if hasattr(sheet, "_images"):
            for img in sheet._images:
                try:
                    image_bytes = img._data()
                    if len(image_bytes) < MIN_IMAGE_SIZE:
                        continue

                    # 确定扩展名
                    ext = img.format.lower() if img.format else "png"
                    if ext == "jpeg":
                        ext = "jpg"

                    image_count += 1
                    image_filename = f"xlsx_{sheet_name}_img{image_count}.{ext}"
                    # 清理文件名中的非法字符
                    image_filename = "".join(c if c.isalnum() or c in "._-" else "_" for c in image_filename)
                    image_path = os.path.join(output_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                except Exception:
                    continue

    return image_count


# -------------------- 核心数据 --------------------

SUPPORTED_FORMATS = {
    ".pptx": extract_images_from_pptx,
    ".docx": extract_images_from_docx,
    ".xlsx": extract_images_from_xlsx,
}


def scan_directory(directory: str) -> List[str]:
    """扫描目录中所有支持的文档文件（不包含子目录）。"""
    files = []
    for entry in os.listdir(directory):
        path = os.path.join(directory, entry)
        if os.path.isfile(path) and Path(path).suffix.lower() in SUPPORTED_FORMATS:
            files.append(path)
    return files


def extract_images_from_file(filepath: str, output_dir: str) -> int:
    """从单个文件中提取图片。"""
    filepath = os.path.abspath(filepath)
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"文件不存在: {filepath}")

    ext = Path(filepath).suffix.lower()
    if ext not in SUPPORTED_FORMATS:
        raise ValueError(
            f"不支持的文件格式: {ext}\n"
            f"支持的格式: {', '.join(SUPPORTED_FORMATS.keys())}"
        )

    os.makedirs(output_dir, exist_ok=True)
    extractor = SUPPORTED_FORMATS[ext]
    return extractor(filepath, output_dir)


def extract_images(
    inputs: List[str],
    output_dir: str,
    flat: bool = False
) -> Dict[str, int]:
    """
    批量提取图片。

    参数:
        inputs: 文件路径或目录路径列表
        output_dir: 根输出目录
        flat: 为 True 时所有图片输出到同一目录；为 False 时按文件名分子目录

    返回:
        {文件路径: 提取图片数量}
    """
    results = {}

    # 收集所有文件
    all_files = []
    for item in inputs:
        item = os.path.abspath(item)
        if os.path.isdir(item):
            files = scan_directory(item)
            all_files.extend(files)
            if not files:
                print(f"⚠️ 目录中未找到支持的文档: {item}")
        elif os.path.isfile(item):
            all_files.append(item)
        else:
            print(f"⚠️ 路径不存在或不是文件/目录: {item}")

    if not all_files:
        print("⚠️ 没有找到可处理的文件")
        return results

    # 去重并保持顺序
    seen = set()
    unique_files = [f for f in all_files if not (f in seen or seen.add(f))]

    total_files = len(unique_files)
    print(f"📦 共发现 {total_files} 个待处理文件")
    print("=" * 50)

    # 如果输出目录不存在，提示正在创建
    output_dir = os.path.abspath(output_dir)
    if not os.path.exists(output_dir):
        print(f"📁 输出目录不存在，正在创建: {output_dir}")
    os.makedirs(output_dir, exist_ok=True)
    total_images = 0

    for idx, filepath in enumerate(unique_files, start=1):
        basename = Path(filepath).stem
        ext = Path(filepath).suffix.lower()

        if flat:
            # 所有图片放一起，文件名加前缀区分
            file_output_dir = output_dir
            prefix = f"{basename}_"
        else:
            # 每个文件一个子目录
            file_output_dir = os.path.join(output_dir, basename)
            prefix = ""

        print(f"\n[{idx}/{total_files}] 📄 {filepath}")
        print(f"     输出到: {file_output_dir}")

        try:
            # 创建临时子目录（flat 模式下不需要）
            if not flat:
                os.makedirs(file_output_dir, exist_ok=True)

            # 提取到临时目录，再重命名（flat 模式下需要加前缀）
            import tempfile
            tmp_dir = tempfile.mkdtemp()
            count = extract_images_from_file(filepath, tmp_dir)

            if flat and count > 0:
                # flat 模式：把文件移到根目录并加前缀
                for fname in os.listdir(tmp_dir):
                    src = os.path.join(tmp_dir, fname)
                    dst = os.path.join(file_output_dir, f"{prefix}{fname}")
                    # 处理重名
                    counter = 1
                    orig_dst = dst
                    while os.path.exists(dst):
                        stem = Path(orig_dst).stem
                        suffix = Path(orig_dst).suffix
                        dst = os.path.join(file_output_dir, f"{stem}_{counter}{suffix}")
                        counter += 1
                    os.rename(src, dst)
                os.rmdir(tmp_dir)
            elif count > 0:
                # 子目录模式：直接移过去
                for fname in os.listdir(tmp_dir):
                    src = os.path.join(tmp_dir, fname)
                    dst = os.path.join(file_output_dir, fname)
                    counter = 1
                    orig_dst = dst
                    while os.path.exists(dst):
                        stem = Path(orig_dst).stem
                        suffix = Path(orig_dst).suffix
                        dst = os.path.join(file_output_dir, f"{stem}_{counter}{suffix}")
                        counter += 1
                    os.rename(src, dst)
                os.rmdir(tmp_dir)
            else:
                # 没提取到图片，清理临时目录
                import shutil
                shutil.rmtree(tmp_dir, ignore_errors=True)
                if not flat:
                    # 空目录也清理
                    try:
                        os.rmdir(file_output_dir)
                    except OSError:
                        pass

            results[filepath] = count
            total_images += count
            print(f"     ✅ 提取 {count} 张图片")

        except Exception as e:
            results[filepath] = 0
            print(f"     ❌ 失败: {e}")

    # 汇总
    print("\n" + "=" * 50)
    print(f"🎉 处理完成: {total_files} 个文件, 共 {total_images} 张图片")
    print(f"📁 输出目录: {os.path.abspath(output_dir)}")
    return results


# -------------------- CLI --------------------

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="从文档中提取内嵌图片",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  %(prog)s report.docx
  %(prog)s report.docx slides.pptx -o ./images
  %(prog)s --dir ./docs -o ./images
  %(prog)s --dir ./docs --flat -o ./all_images
        """,
    )
    parser.add_argument(
        "inputs",
        nargs="*",
        help="要处理的文件或目录路径（可多选）",
    )
    parser.add_argument(
        "-d", "--dir",
        action="append",
        default=[],
        help="扫描目录中的所有支持文件（可多次使用）",
    )
    parser.add_argument(
        "-o", "--output",
        default="extracted_images",
        help="输出目录（默认: extracted_images）",
    )
    parser.add_argument(
        "--flat",
        action="store_true",
        help="扁平输出：所有图片放到同一目录，不按文件分目录",
    )
    return parser


def main():
    parser = build_parser()
    args = parser.parse_args()

    # 收集所有输入
    all_inputs = list(args.inputs) if args.inputs else []
    for d in args.dir:
        all_inputs.append(d)

    if not all_inputs:
        parser.print_help()
        sys.exit(1)

    extract_images(all_inputs, args.output, flat=args.flat)


if __name__ == "__main__":
    main()
